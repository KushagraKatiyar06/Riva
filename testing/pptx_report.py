"""
testing/pptx_report.py
Generates a brand-aware competitive intelligence PowerPoint deck.
Each company's section uses their extracted brand color as accent.

Requires:
    pip install python-pptx Pillow
"""

import os
import io
import re
import time
import base64
import json
import requests
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from report import generate_intel, scrape_brand_assets, get_available_domains, chunks_for

load_dotenv()

GEMINI_KEY  = os.getenv("GEMINI_API_KEY")
REPORTS_DIR = Path(__file__).parent / "reports"
REPORTS_DIR.mkdir(exist_ok=True)


DARK_BG    = RGBColor(0x07, 0x0D, 0x18)
MID_BG     = RGBColor(0x0C, 0x18, 0x2C)
CARD_BG    = RGBColor(0x0F, 0x1E, 0x36)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xC8, 0xD8, 0xE8)
DIM_TEXT   = RGBColor(0x8A, 0xA4, 0xBC)
RED        = RGBColor(0xE0, 0x50, 0x50)

# Fallback accent colors when brand color not found
_FALLBACK_RGB = [
    (0x00, 0xCC, 0xCC),  # cyan
    (0xF3, 0x81, 0x0F),  # orange
    (0x76, 0x45, 0xFF),  # purple
    (0x00, 0xC8, 0x5A),  # green
]

W = Inches(13.333)
H = Inches(7.5)


def _parse_hex(h: str):
    """Return (r, g, b) from '#RRGGBB' or '#RGB', or None on failure."""
    if not h:
        return None
    try:
        h = h.lstrip("#")
        if len(h) == 3:
            h = h[0] * 2 + h[1] * 2 + h[2] * 2
        if len(h) == 6:
            return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        pass
    return None


def _make_theme(images: dict, domains: list) -> dict:
    """
    Build per-domain color theme from brand_color in images dict.
    Returns {domain: {"accent": RGBColor, "dark": RGBColor, "r": int, "g": int, "b": int}}
    """
    theme = {}
    for i, d in enumerate(domains):
        hex_color = (images.get(d) or {}).get("brand_color")
        rgb = _parse_hex(hex_color) if hex_color else None
        if rgb is None:
            rgb = _FALLBACK_RGB[i % len(_FALLBACK_RGB)]
        r, g, b = rgb
        # Dark version: 18% brightness — used as card backgrounds
        dr, dg, db = max(0, int(r * 0.18)), max(0, int(g * 0.18)), max(0, int(b * 0.18))
        theme[d] = {
            "accent": RGBColor(r, g, b),
            "dark":   RGBColor(dr, dg, db),
            "r": r, "g": g, "b": b,
        }
    return theme


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor = None, line: RGBColor = None, line_w: int = 0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_w:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def _truncate(text: str, max_chars: int) -> str:
    text = str(text).strip()
    return text[:max_chars - 1] + "…" if len(text) > max_chars else text


def add_text(slide, text: str, x, y, w, h,
             size: int = 16, bold: bool = False, color: RGBColor = WHITE,
             align=PP_ALIGN.LEFT, wrap: bool = True, max_chars: int = 300) -> None:
    text = _truncate(str(text), max_chars)
    # Auto-shrink font for long strings to prevent overflow
    if len(text) > 220 and size > 12:
        size = max(10, size - 4)
    elif len(text) > 130 and size > 14:
        size = max(11, size - 2)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color


def add_text_lines(slide, lines: list, x, y, w, h,
                   size: int = 13, color: RGBColor = LIGHT_TEXT,
                   bullet: bool = True, max_items: int = 5) -> None:
    # Filter and cap
    lines = [l for l in lines if l][:max_items]
    if not lines:
        return
    # Adaptive sizing
    n = len(lines)
    if n >= 5:
        size = max(10, size - 1)
    if n >= 6:
        size = max(9, size - 1)

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        run = p.add_run()
        line_str = _truncate(str(line) if not isinstance(line, dict)
                             else " — ".join(str(v) for v in line.values() if v), 140)
        run.text = ("• " if bullet else "") + line_str
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def add_logo(slide, b64_uri: str, x, y, w=Inches(0.6), h=Inches(0.6)):
    if not b64_uri or "base64," not in b64_uri:
        return
    try:
        data = base64.b64decode(b64_uri.split("base64,")[1])
        slide.shapes.add_picture(io.BytesIO(data), x, y, width=w, height=h)
    except Exception:
        pass


def slide_title(prs, intel: dict, images: dict, theme: dict = None):
    """Split-panel title slide using each company's brand color."""
    domains = list(intel["domains"].keys())
    if theme is None:
        theme = _make_theme(images, domains)

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    half = W / 2

    # Left panel — domain A color
    if domains:
        t = theme[domains[0]]
        add_rect(slide, Inches(0), Inches(0), half, H, fill=t["dark"])
        add_rect(slide, Inches(0), Inches(0), Inches(0.06), H, fill=t["accent"])

    # Right panel — domain B color
    if len(domains) > 1:
        t = theme[domains[1]]
        add_rect(slide, half, Inches(0), half, H, fill=t["dark"])
        add_rect(slide, W - Inches(0.06), Inches(0), Inches(0.06), H, fill=t["accent"])

    # Center divider line
    add_rect(slide, half - Inches(0.01), Inches(0), Inches(0.02), H, fill=DIM_TEXT)

    # RIVA label
    add_text(slide, "RIVA COMPETITIVE INTELLIGENCE",
             Inches(0.5), Inches(0.28), W - Inches(1), Inches(0.4),
             size=9, bold=True, color=DIM_TEXT, align=PP_ALIGN.CENTER)

    # Main title
    add_text(slide, "Competitive Analysis",
             Inches(0.5), Inches(0.72), W - Inches(1), Inches(0.85),
             size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER, max_chars=80)

    # Domain names line
    names = " vs ".join(intel["domains"][d]["display_name"] for d in domains)
    # Use a gradient-like effect: write each name in its color
    if len(domains) == 2:
        d0, d1 = domains[0], domains[1]
        n0 = intel["domains"][d0]["display_name"]
        n1 = intel["domains"][d1]["display_name"]
        add_text(slide, n0 + "  ×  " + n1,
                 Inches(0.5), Inches(1.55), W - Inches(1), Inches(0.55),
                 size=20, color=LIGHT_TEXT, align=PP_ALIGN.CENTER, max_chars=80)
    else:
        add_text(slide, names, Inches(0.5), Inches(1.55), W - Inches(1), Inches(0.55),
                 size=20, color=LIGHT_TEXT, align=PP_ALIGN.CENTER, max_chars=80)

    # Per-domain info cards (bottom half)
    x_positions = [Inches(0.5), half + Inches(0.5)]
    card_w = half - Inches(1.0)

    for i, d in enumerate(domains):
        info = intel["domains"][d]
        t    = theme[d]
        x    = x_positions[i]

        # Logo
        logo = (images.get(d) or {}).get("logo")
        add_logo(slide, logo, x, Inches(2.55), w=Inches(0.75), h=Inches(0.75))

        # Name + tagline
        add_text(slide, info["display_name"],
                 x + Inches(0.9), Inches(2.65), card_w - Inches(0.95), Inches(0.45),
                 size=18, bold=True, color=WHITE, max_chars=40)
        add_text(slide, info["tagline"],
                 x + Inches(0.9), Inches(3.1), card_w - Inches(0.95), Inches(0.55),
                 size=11, color=LIGHT_TEXT, max_chars=100)

        # Target audience badge
        add_rect(slide, x, Inches(3.8), card_w, Inches(0.38), fill=t["dark"])
        add_rect(slide, x, Inches(3.8), Inches(0.04), Inches(0.38), fill=t["accent"])
        add_text(slide, info.get("target_audience", ""),
                 x + Inches(0.12), Inches(3.85), card_w - Inches(0.2), Inches(0.32),
                 size=11, color=LIGHT_TEXT, max_chars=80)

        # Top 3 features
        add_text(slide, "KEY CAPABILITIES",
                 x, Inches(4.32), card_w, Inches(0.28),
                 size=8, bold=True, color=t["accent"])
        add_text_lines(slide, info.get("top_features", [])[:3],
                       x, Inches(4.6), card_w, Inches(1.1),
                       size=11, color=LIGHT_TEXT, max_items=3)

    # Date + footer
    add_text(slide, datetime.now().strftime("%B %d, %Y"),
             Inches(0), Inches(7.1), W, Inches(0.32),
             size=9, color=DIM_TEXT, align=PP_ALIGN.CENTER)


def slide_exec_summary(prs, intel: dict, images: dict, theme: dict = None):
    domains = list(intel["domains"].keys())
    if theme is None:
        theme = _make_theme(images, domains)

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    # Top label bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.5), fill=MID_BG)
    add_text(slide, "EXECUTIVE SUMMARY",
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.35),
             size=10, bold=True, color=DIM_TEXT)

    col_w = Inches(6.1)
    x_positions = [Inches(0.3), Inches(6.9)]

    for i, d in enumerate(domains):
        info = intel["domains"][d]
        t    = theme[d]
        x    = x_positions[i]

        # Card background
        add_rect(slide, x, Inches(0.58), col_w, Inches(6.7), fill=CARD_BG)
        # Brand color stripe at top of card
        add_rect(slide, x, Inches(0.58), col_w, Inches(0.05), fill=t["accent"])

        # Logo + name
        logo = (images.get(d) or {}).get("logo")
        add_logo(slide, logo, x + Inches(0.2), Inches(0.75), w=Inches(0.55), h=Inches(0.55))
        add_text(slide, info["display_name"],
                 x + Inches(0.88), Inches(0.8), col_w - Inches(1.1), Inches(0.42),
                 size=18, bold=True, color=WHITE, max_chars=35)
        add_text(slide, info["target_audience"],
                 x + Inches(0.2), Inches(1.38), col_w - Inches(0.4), Inches(0.35),
                 size=10, color=t["accent"], max_chars=90)

        # Tagline
        add_text(slide, info["tagline"],
                 x + Inches(0.2), Inches(1.75), col_w - Inches(0.4), Inches(0.7),
                 size=12, color=LIGHT_TEXT, max_chars=130)

        # Features section
        add_text(slide, "TOP FEATURES",
                 x + Inches(0.2), Inches(2.55), col_w - Inches(0.4), Inches(0.28),
                 size=8, bold=True, color=DIM_TEXT)
        add_text_lines(slide, info.get("top_features", [])[:5],
                       x + Inches(0.2), Inches(2.85), col_w - Inches(0.4), Inches(1.85),
                       size=11, color=LIGHT_TEXT, max_items=5)

        # Strengths
        add_text(slide, "STRENGTHS",
                 x + Inches(0.2), Inches(4.82), col_w - Inches(0.4), Inches(0.28),
                 size=8, bold=True, color=t["accent"])
        add_text_lines(slide, info.get("strengths", [])[:3],
                       x + Inches(0.2), Inches(5.12), col_w - Inches(0.4), Inches(1.0),
                       size=11, color=LIGHT_TEXT, max_items=3)

        # Weaknesses
        add_text(slide, "GAPS",
                 x + Inches(0.2), Inches(6.22), col_w - Inches(0.4), Inches(0.28),
                 size=8, bold=True, color=RED)
        add_text_lines(slide, info.get("weaknesses", [])[:2],
                       x + Inches(0.2), Inches(6.5), col_w - Inches(0.4), Inches(0.65),
                       size=10, color=LIGHT_TEXT, bullet=True, max_items=2)


def slide_pricing(prs, intel: dict, theme: dict = None, images: dict = None):
    domains = list(intel["domains"].keys())
    if theme is None:
        theme = _make_theme(images or {}, domains)

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.5), fill=MID_BG)
    add_text(slide, "PRICING COMPARISON",
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.35),
             size=10, bold=True, color=DIM_TEXT)

    col_w = Inches(6.1)
    x_positions = [Inches(0.3), Inches(6.9)]

    for i, d in enumerate(domains):
        info  = intel["domains"][d]
        t     = theme[d]
        x     = x_positions[i]
        tiers = info.get("pricing_tiers", [])[:4]

        add_text(slide, info["display_name"],
                 x, Inches(0.58), col_w, Inches(0.45),
                 size=16, bold=True, color=WHITE, max_chars=40)
        add_rect(slide, x, Inches(0.95), col_w - Inches(0.1), Inches(0.03), fill=t["accent"])

        y = Inches(1.05)
        tier_h = Inches(1.52)
        for tier in tiers:
            add_rect(slide, x, y, col_w - Inches(0.15), tier_h, fill=CARD_BG)
            add_rect(slide, x, y, Inches(0.04), tier_h, fill=t["accent"])

            tier_name = _truncate(str(tier.get("name", "")), 30)
            tier_price = _truncate(str(tier.get("price", "")), 25)
            add_text(slide, tier_name.upper(),
                     x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.3),
                     size=9, bold=True, color=t["accent"])
            add_text(slide, tier_price,
                     x + Inches(0.15), y + Inches(0.36), Inches(2.8), Inches(0.45),
                     size=17, bold=True, color=WHITE, max_chars=25)

            highlights = [str(h) for h in (tier.get("highlights") or [])[:3]]
            add_text_lines(slide, highlights,
                           x + Inches(0.15), y + Inches(0.82), col_w - Inches(0.4), Inches(0.72),
                           size=10, color=LIGHT_TEXT, max_items=3)
            y += tier_h + Inches(0.1)


def slide_differentiators(prs, intel: dict, theme: dict = None, images: dict = None):
    domains = list(intel["domains"].keys())
    if theme is None:
        theme = _make_theme(images or {}, domains)

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.5), fill=MID_BG)
    add_text(slide, "COMPETITIVE DIFFERENTIATORS",
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.35),
             size=10, bold=True, color=DIM_TEXT)

    comparison  = intel["comparison"]
    col_w       = Inches(6.1)
    x_positions = [Inches(0.3), Inches(6.9)]

    for i, d in enumerate(domains):
        info  = intel["domains"][d]
        t     = theme[d]
        x     = x_positions[i]
        key   = "only_in_" + d.replace(".", "_")
        items = comparison.get(key, [])[:7]

        add_text(slide, "ONLY IN " + info["display_name"].upper(),
                 x, Inches(0.58), col_w, Inches(0.38),
                 size=13, bold=True, color=t["accent"], max_chars=50)

        y = Inches(1.05)
        for item in items:
            add_rect(slide, x, y, col_w - Inches(0.2), Inches(0.58), fill=CARD_BG)
            add_rect(slide, x, y, Inches(0.04), Inches(0.58), fill=t["accent"])
            add_text(slide, item,
                     x + Inches(0.15), y + Inches(0.08),
                     col_w - Inches(0.45), Inches(0.46),
                     size=10, color=WHITE, max_chars=100)
            y += Inches(0.64)

        # Weaknesses at bottom
        add_text(slide, "GAPS TO ADDRESS",
                 x, Inches(5.62), col_w, Inches(0.28),
                 size=8, bold=True, color=RED)
        add_text_lines(slide, info.get("weaknesses", [])[:3],
                       x, Inches(5.92), col_w, Inches(1.3),
                       size=10, color=LIGHT_TEXT, max_items=3)


def slide_gtm(prs, intel: dict, gtm: dict, theme: dict = None, images: dict = None):
    domains = list(intel["domains"].keys())
    if theme is None:
        theme = _make_theme(images or {}, domains)
    t0 = theme[domains[0]] if domains else {"accent": RGBColor(0, 204, 204), "dark": CARD_BG}
    t1 = theme[domains[1]] if len(domains) > 1 else {"accent": RGBColor(0, 200, 90), "dark": CARD_BG}

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.5), fill=MID_BG)
    add_text(slide, "GTM STRATEGY & RECOMMENDATION",
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.35),
             size=10, bold=True, color=DIM_TEXT)

    comparison = intel["comparison"]

    # Pricing verdict card
    add_rect(slide, Inches(0.3), Inches(0.6), Inches(6.1), Inches(1.5), fill=CARD_BG)
    add_rect(slide, Inches(0.3), Inches(0.6), Inches(0.04), Inches(1.5), fill=t0["accent"])
    add_text(slide, "PRICING VERDICT",
             Inches(0.5), Inches(0.68), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=DIM_TEXT)
    add_text(slide, comparison.get("pricing_verdict", ""),
             Inches(0.5), Inches(0.96), Inches(5.8), Inches(1.0),
             size=12, color=WHITE, max_chars=260)

    # Positioning verdict card
    add_rect(slide, Inches(6.9), Inches(0.6), Inches(6.1), Inches(1.5), fill=CARD_BG)
    add_rect(slide, Inches(6.9), Inches(0.6), Inches(0.04), Inches(1.5), fill=t1["accent"])
    add_text(slide, "POSITIONING VERDICT",
             Inches(7.1), Inches(0.68), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=DIM_TEXT)
    add_text(slide, comparison.get("positioning_verdict", ""),
             Inches(7.1), Inches(0.96), Inches(5.8), Inches(1.0),
             size=12, color=WHITE, max_chars=260)

    # Key messages
    add_text(slide, "KEY MESSAGES",
             Inches(0.5), Inches(2.28), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=DIM_TEXT)
    add_text_lines(slide, gtm.get("key_messages", []),
                   Inches(0.3), Inches(2.58), Inches(6.1), Inches(2.0),
                   size=12, color=LIGHT_TEXT, max_items=4)

    # Objections
    add_text(slide, "OBJECTION HANDLING",
             Inches(7.1), Inches(2.28), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=DIM_TEXT)
    add_text_lines(slide, gtm.get("objections", []),
                   Inches(6.9), Inches(2.58), Inches(6.1), Inches(2.0),
                   size=12, color=LIGHT_TEXT, max_items=4)

    # Recommendation banner
    add_rect(slide, Inches(0), Inches(4.75), W, Inches(2.55), fill=MID_BG)
    add_rect(slide, Inches(0), Inches(4.75), W, Inches(0.04), fill=t0["accent"])
    add_text(slide, "RECOMMENDATION",
             Inches(0.5), Inches(4.9), Inches(12), Inches(0.3),
             size=8, bold=True, color=t0["accent"])
    add_text(slide, comparison.get("recommendation", ""),
             Inches(0.5), Inches(5.22), Inches(12.3), Inches(1.9),
             size=13, color=WHITE, max_chars=480)


def slide_action_items(prs, gtm: dict, theme: dict = None, domains: list = None):
    if theme is None or domains is None:
        theme = {}
        domains = []
    t0 = theme.get(domains[0], {"accent": RGBColor(0, 204, 204)}) if domains else {"accent": RGBColor(0, 204, 204)}
    t1 = theme.get(domains[1], {"accent": RGBColor(0, 200, 90)}) if len(domains) > 1 else {"accent": RGBColor(0, 200, 90)}

    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.5), fill=MID_BG)
    add_text(slide, "ACTION ITEMS",
             Inches(0.5), Inches(0.1), Inches(10), Inches(0.35),
             size=10, bold=True, color=DIM_TEXT)

    actions = [a for a in gtm.get("action_items", []) if a][:8]
    mid     = (len(actions) + 1) // 2
    cols    = [actions[:mid], actions[mid:]]
    x_positions = [Inches(0.3), Inches(6.9)]
    accents = [t0["accent"], t1["accent"]]

    for col_items, x, accent in zip(cols, x_positions, accents):
        y = Inches(0.65)
        offset = 0 if x == x_positions[0] else mid
        for j, action in enumerate(col_items):
            item_h = Inches(0.92)
            add_rect(slide, x, y, Inches(6.1), item_h, fill=CARD_BG)
            add_rect(slide, x, y, Inches(0.48), item_h, fill=accent)
            add_text(slide, str(j + 1 + offset),
                     x + Inches(0.08), y + Inches(0.22), Inches(0.34), Inches(0.48),
                     size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
            add_text(slide, action,
                     x + Inches(0.58), y + Inches(0.14), Inches(5.4), Inches(0.68),
                     size=11, color=WHITE, max_chars=150)
            y += item_h + Inches(0.08)

    add_text(slide, "RIVA COMPETITIVE INTELLIGENCE  ·  CONFIDENTIAL",
             Inches(0), Inches(7.18), W, Inches(0.28),
             size=8, color=DIM_TEXT, align=PP_ALIGN.CENTER)


def generate_gtm(domains: list, intel: dict, focus: str = None) -> dict:
    client = genai.Client(api_key=GEMINI_KEY)

    context = ""
    for d in domains:
        context += (
            f"\n=== {d} ===\n"
            f"{chunks_for(d, 'pricing features positioning target audience competitors')}\n"
        )

    names = [intel["domains"][d]["display_name"] for d in domains]
    focus_instruction = (
        f"\n\nFOCUS: Tailor key_messages, objections, and action_items specifically around **{focus}**."
    ) if focus else ""

    prompt = (
        "You are a B2B marketing strategist. Based on the competitive context below, "
        "generate actionable GTM content.\n\n"
        "CONTEXT:\n" + context + "\n\n"
        "Companies: " + " vs ".join(names) + "\n\n"
        "Return ONLY valid JSON:\n"
        "{\n"
        '  "key_messages": ["3-4 key selling messages"],\n'
        '  "objections": ["3-4 objections and one-line rebuttals"],\n'
        '  "action_items": ["6-8 specific next steps"]\n'
        "}" + focus_instruction
    )

    for attempt in range(3):
        res  = client.models.generate_content(model="gemini-2.5-flash", contents=prompt)
        text = res.text.strip()
        if "```" in text:
            text = text.split("```")[1].lstrip("json").strip()
            text = text.split("```")[0].strip()
        start, end = text.find('{'), text.rfind('}')
        if start != -1 and end != -1:
            text = text[start:end + 1]
        try:
            result = json.loads(text)
            def _to_str(item):
                if isinstance(item, str): return item
                if isinstance(item, dict): return " — ".join(str(v) for v in item.values() if v)
                return str(item)
            for key in ("key_messages", "objections", "action_items"):
                if key in result and isinstance(result[key], list):
                    result[key] = [_to_str(i) for i in result[key]]
            return result
        except json.JSONDecodeError as e:
            print(f"  generate_gtm JSON parse error (attempt {attempt+1}/3): {e}")
            if attempt == 2:
                raise
            time.sleep(5)
    raise RuntimeError("generate_gtm failed after 3 attempts")


def render_preview_html(intel: dict, gtm: dict, images: dict, theme: dict) -> str:
    """Generate a slide-deck HTML preview of the PPTX content — shown in the frontend iframe."""
    domains    = list(intel["domains"].keys())
    comparison = intel["comparison"]
    date_str   = datetime.now().strftime("%B %d, %Y")

    def _hex(d: str) -> str:
        t = theme.get(d, {})
        acc = t.get("accent")
        if acc is not None:
            return f"#{str(acc)}"  # str(RGBColor) → 'RRGGBB' uppercase hex
        return (images.get(d) or {}).get("brand_color") or "#4db8ff"

    def _dark_hex(d: str) -> str:
        t = theme.get(d, {})
        dk = t.get("dark")
        if dk is not None:
            return f"#{str(dk)}"
        return "#0f1e36"

    colors = {d: _hex(d)      for d in domains}
    darks  = {d: _dark_hex(d) for d in domains}
    c0     = colors.get(domains[0], "#4db8ff") if domains          else "#4db8ff"
    c1     = colors.get(domains[1], "#ff6b6b") if len(domains) > 1 else "#ff6b6b"
    n0     = intel["domains"][domains[0]]["display_name"] if domains          else ""
    n1     = intel["domains"][domains[1]]["display_name"] if len(domains) > 1 else ""

    def _li(items):
        return "".join(f"<li>{str(i)}</li>" for i in items if i)

    d0 = intel["domains"].get(domains[0], {}) if domains          else {}
    d1 = intel["domains"].get(domains[1], {}) if len(domains) > 1 else {}

    # Slide 1 — Title
    s1 = f'''<div class="slide s-title">
      <div class="half" style="background:{darks.get(domains[0] if domains else "","#0f1e36")};border-right:2px solid {c0}44;">
        <div class="brand-bar" style="background:{c0};"></div>
        <div class="hn" style="color:{c0};">{n0}</div>
        <div class="ht">{d0.get("tagline","")[:130]}</div>
      </div>
      <div class="half" style="background:{darks.get(domains[1] if len(domains)>1 else "","#0f1e36")};">
        <div class="brand-bar" style="background:{c1};"></div>
        <div class="hn" style="color:{c1};">{n1}</div>
        <div class="ht">{d1.get("tagline","")[:130]}</div>
      </div>
      <div class="tc">
        <div class="tl">RIVA COMPETITIVE INTELLIGENCE</div>
        <div class="tm">Competitive Analysis</div>
        <div class="tv">{n0} &times; {n1}</div>
        <div class="td">{date_str}</div>
      </div>
    </div>'''

    # Slide 2 — Executive Summary
    ec = ""
    for d in domains:
        info = intel["domains"][d]; c = colors[d]
        ec += f'''<div class="ec" style="border-top:3px solid {c};">
          <div class="cn" style="color:{c};">{info["display_name"]}</div>
          <div class="cs">{info.get("target_audience","")}</div>
          <div class="lbl">TOP FEATURES</div><ul class="il">{_li(info.get("top_features",[])[:5])}</ul>
          <div class="lbl" style="color:#6dc08a;">STRENGTHS</div><ul class="il">{_li(info.get("strengths",[])[:3])}</ul>
          <div class="lbl" style="color:#e05555;">GAPS</div><ul class="il">{_li(info.get("weaknesses",[])[:2])}</ul>
        </div>'''
    s2 = f'<div class="slide"><div class="sh">Executive Summary</div><div class="tc2">{ec}</div></div>'

    # Slide 3 — Pricing
    pc = ""
    for d in domains:
        info = intel["domains"][d]; c = colors[d]
        th = ""
        for tier in info.get("pricing_tiers", [])[:5]:
            hl = _li((tier.get("highlights") or [])[:3])
            th += f'<div class="tier" style="border-left:3px solid {c}55;"><div class="tnm" style="color:{c};">{tier.get("name","")}</div><div class="tpr">{tier.get("price","")}</div><ul class="ti">{hl}</ul></div>'
        pc += f'<div><div class="cn" style="color:{c};">{info["display_name"]}</div>{th}</div>'
    s3 = f'<div class="slide"><div class="sh">Pricing Comparison</div><div class="tc2">{pc}</div></div>'

    # Slide 4 — Differentiators
    dc = ""
    for d in domains:
        info = intel["domains"][d]; c = colors[d]
        key  = "only_in_" + d.replace(".", "_")
        dc += f'''<div>
          <div class="cn" style="color:{c};">Only in {info["display_name"]}</div>
          <ul class="il">{_li(comparison.get(key,[])[:7])}</ul>
          <div class="lbl" style="color:#e05555;margin-top:12px;">GAPS</div>
          <ul class="il">{_li(info.get("weaknesses",[])[:3])}</ul>
        </div>'''
    s4 = f'<div class="slide"><div class="sh">Competitive Differentiators</div><div class="tc2">{dc}</div></div>'

    # Slide 5 — GTM
    km = _li(gtm.get("key_messages", []))
    ob = _li(gtm.get("objections", []))
    s5 = f'''<div class="slide"><div class="sh">GTM Strategy &amp; Recommendation</div>
      <div class="tc2">
        <div class="vbox" style="border-left:4px solid {c0};"><div class="lbl">PRICING VERDICT</div><div class="vt">{comparison.get("pricing_verdict","")}</div></div>
        <div class="vbox" style="border-left:4px solid {c1};"><div class="lbl">POSITIONING VERDICT</div><div class="vt">{comparison.get("positioning_verdict","")}</div></div>
      </div>
      <div class="tc2" style="margin-top:14px;">
        <div><div class="lbl">KEY MESSAGES</div><ul class="il">{km}</ul></div>
        <div><div class="lbl">OBJECTION HANDLING</div><ul class="il">{ob}</ul></div>
      </div>
      <div class="rec" style="border-top:3px solid {c0};margin-top:16px;">
        <div class="lbl" style="color:{c0};">RECOMMENDATION</div>
        <div class="rt">{comparison.get("recommendation","")}</div>
      </div>
    </div>'''

    # Slide 6 — Action Items
    action_items = [item for item in gtm.get("action_items", []) if item]
    ah = "".join(
        f'<div class="ai"><div class="an" style="background:{c0};">{i+1}</div><div class="at">{item}</div></div>'
        for i, item in enumerate(action_items)
    )
    s6 = f'<div class="slide"><div class="sh">Action Items</div><div class="ag">{ah}</div></div>'

    slides = [("Title", s1), ("Executive Summary", s2), ("Pricing Comparison", s3),
              ("Competitive Differentiators", s4), ("GTM Strategy", s5), ("Action Items", s6)]
    body = "\n".join(
        f'<div class="sw"><div class="sn">Slide {i+1} — {nm}</div>{html}</div>'
        for i, (nm, html) in enumerate(slides)
    )

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<base target="_blank">
<title>Deck Preview — {n0} vs {n1}</title>
<style>
  *,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
  body{{background:#050912;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;font-size:13px;color:#c8d8e8;padding:24px 0 40px}}
  .page{{max-width:1100px;margin:0 auto;padding:0 20px}}
  .ph{{margin-bottom:24px;padding-bottom:12px;border-bottom:1px solid rgba(255,255,255,0.08)}}
  .ph-t{{font-size:10px;font-weight:700;letter-spacing:3px;color:rgba(255,255,255,0.3);text-transform:uppercase}}
  .ph-s{{font-size:15px;font-weight:600;color:white;margin-top:4px}}
  .sw{{margin-bottom:24px}}
  .sn{{font-size:9px;font-weight:700;letter-spacing:2px;color:rgba(255,255,255,0.2);text-transform:uppercase;margin-bottom:5px}}
  .slide{{background:#070d18;border:1px solid rgba(255,255,255,0.07);border-radius:10px;padding:26px 30px;overflow:hidden}}
  .sh{{font-size:16px;font-weight:700;color:white;margin-bottom:16px;padding-bottom:10px;border-bottom:1px solid rgba(255,255,255,0.07)}}
  .tc2{{display:grid;grid-template-columns:1fr 1fr;gap:22px}}
  .lbl{{font-size:9px;font-weight:700;letter-spacing:2px;text-transform:uppercase;color:rgba(255,255,255,0.32);margin:10px 0 4px}}
  .il{{padding-left:14px;list-style:disc}}.il li{{margin-bottom:4px;line-height:1.5;color:#c8d8e8;font-size:12px}}
  .s-title{{position:relative;min-height:200px;padding:0;display:flex}}
  .half{{flex:1;padding:26px 22px;position:relative}}
  .brand-bar{{position:absolute;top:0;left:0;right:0;height:3px}}
  .hn{{font-size:22px;font-weight:700;margin-top:10px;margin-bottom:6px}}
  .ht{{font-size:11.5px;color:rgba(255,255,255,0.45);line-height:1.5}}
  .tc{{position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);text-align:center;pointer-events:none}}
  .tl{{font-size:9px;font-weight:700;letter-spacing:3px;color:rgba(255,255,255,0.28);text-transform:uppercase;margin-bottom:5px}}
  .tm{{font-size:24px;font-weight:800;color:white}}
  .tv{{font-size:13px;color:rgba(255,255,255,0.45);margin-top:5px}}
  .td{{font-size:10px;color:rgba(255,255,255,0.22);margin-top:6px}}
  .ec{{padding-top:10px}}.cn{{font-size:15px;font-weight:700;margin-bottom:2px}}.cs{{font-size:11px;color:rgba(255,255,255,0.4);margin-bottom:6px}}
  .tier{{background:rgba(255,255,255,0.03);border-radius:6px;padding:9px 11px;margin-bottom:7px}}
  .tnm{{font-size:9px;font-weight:700;text-transform:uppercase;letter-spacing:0.5px;margin-bottom:2px}}
  .tpr{{font-size:15px;font-weight:800;color:white;margin-bottom:4px}}
  .ti{{padding-left:12px;font-size:11px;color:#8aa4bc}}.ti li{{margin-bottom:2px}}
  .vbox{{background:rgba(255,255,255,0.03);border-radius:0 6px 6px 0;padding:11px 13px}}
  .vt{{font-size:12.5px;color:#c8d8e8;line-height:1.6;margin-top:3px}}
  .rec{{background:#0a1628;border-radius:6px;padding:13px 15px}}
  .rt{{font-size:13px;color:#c8d8e8;line-height:1.7;margin-top:5px}}
  .ag{{display:grid;grid-template-columns:1fr 1fr;gap:9px}}
  .ai{{display:flex;align-items:flex-start;gap:9px;background:rgba(255,255,255,0.03);border-radius:6px;padding:9px 11px}}
  .an{{width:22px;height:22px;border-radius:4px;display:flex;align-items:center;justify-content:center;font-size:11px;font-weight:700;color:#070d18;flex-shrink:0}}
  .at{{font-size:12px;color:#c8d8e8;line-height:1.5}}
</style>
</head>
<body>
<div class="page">
  <div class="ph"><div class="ph-t">RIVA COMPETITIVE INTELLIGENCE</div><div class="ph-s">Deck Preview — {n0} vs {n1} &middot; {date_str}</div></div>
  {body}
</div>
</body>
</html>"""


def main():
    missing = [k for k in ("CLOUDFLARE_ACCOUNT_ID", "CLOUDFLARE_API_TOKEN", "GEMINI_API_KEY") if not os.getenv(k)]
    if missing:
        print(f"Missing env vars: {', '.join(missing)}")
        return

    domains = get_available_domains()
    if not domains:
        print("No domains found in testing/extracts/ — run the browser agent first.")
        return

    print("\nRiva PowerPoint Generator")
    print("=" * 40)
    for i, d in enumerate(domains, 1):
        print(f"  {i}. {d}")

    selected = []
    for slot in ["First domain", "Second domain (optional)"]:
        raw = input(f"\n{slot}: ").strip()
        if not raw:
            break
        domain = domains[int(raw) - 1] if raw.isdigit() and 1 <= int(raw) <= len(domains) else raw
        if domain not in domains:
            print(f"  '{domain}' not found.")
            break
        selected.append(domain)

    if not selected:
        print("No domains selected.")
        return

    print(f"\nGenerating deck for: {', '.join(selected)}")
    print("  Scraping brand assets...", end=" ", flush=True)
    images = {d: scrape_brand_assets(d) for d in selected}
    for d, imgs in images.items():
        color = imgs.get("brand_color") or "no color found"
        print(f"\n    {d}: color={color}, logo={'yes' if imgs.get('logo') else 'no'}", end="")
    print()

    theme = _make_theme(images, selected)

    print("  Querying Vectorize + generating intel...", end=" ", flush=True)
    intel = generate_intel(selected)
    print("done")

    print("  Generating GTM strategy...", end=" ", flush=True)
    gtm = generate_gtm(selected, intel)
    print("done")

    print("  Building slides...", end=" ", flush=True)
    prs = new_prs()
    slide_title(prs, intel, images, theme)
    slide_exec_summary(prs, intel, images, theme)
    slide_pricing(prs, intel, theme, images)
    slide_differentiators(prs, intel, theme, images)
    slide_gtm(prs, intel, gtm, theme, images)
    slide_action_items(prs, gtm, theme, selected)
    print("done")

    names    = "_vs_".join(intel["domains"][d]["display_name"].lower().replace(" ", "-") for d in selected)
    ts       = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = REPORTS_DIR / f"riva_deck_{names}_{ts}.pptx"
    prs.save(str(out_path))
    print(f"\n  Deck saved: {out_path}")
    print(f"  {len(prs.slides)} slides — open in PowerPoint or Google Slides.")


if __name__ == "__main__":
    main()
